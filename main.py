import streamlit as st
from groq import Groq
from PyPDF2 import PdfReader
from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration
from docx import Document
from gtts import gTTS
import tempfile
import speech_recognition as sr
from pptx import Presentation/'//////////'
import zipfile

st.set_page_config(page_title="AI Chatbot", page_icon="🤖")

api_key = st.secrets.get("GROQ_API_KEY")

if not api_key:
    st.error("API key missing")
    st.stop()

client = Groq(api_key=api_key)

selected_app = st.sidebar.radio("Choice :", ["ChatBot", "Voice Assistant"])


if selected_app == "ChatBot":
        @st.cache_resource
        def load_model():
            processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
            return processor, model

        processor, model = load_model()


        st.title("🤖 AI Chatbot")

        if "messages" not in st.session_state:
            st.session_state.messages = []

        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])


        uploaded_file = st.file_uploader(
            "Upload file (txt, pdf, image)",
            type=["txt", "pdf", "png", "jpg", "jpeg", "docx", "pptx", "zip"]
        )

        file_text = ""
        image = None

        if uploaded_file:

            if uploaded_file.type == "text/plain":
                file_text = uploaded_file.read().decode("utf-8")

            elif uploaded_file.type == "application/pdf":
                reader = PdfReader(uploaded_file)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        file_text += text

            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(uploaded_file)
                for para in doc.paragraphs:
                    file_text += para.text + "\n"

            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(uploaded_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_text += shape.text + "\n"

           
            elif uploaded_file.type == "application/zip":
                import io

                with zipfile.ZipFile(uploaded_file) as z:
                    for file_name in z.namelist():
                        try:
                            with z.open(file_name) as f:

                                if file_name.endswith(".txt"):
                                    file_text += f.read().decode("utf-8") + "\n"

                                elif file_name.endswith(".pdf"):
                                    reader = PdfReader(io.BytesIO(f.read()))
                                    for page in reader.pages:
                                        text = page.extract_text()
                                        if text:
                                            file_text += text

                                elif file_name.endswith(".docx"):
                                    doc = Document(io.BytesIO(f.read()))
                                    for para in doc.paragraphs:
                                        file_text += para.text + "\n"

                                elif file_name.endswith(".pptx"):
                                    prs = Presentation(io.BytesIO(f.read()))
                                    for slide in prs.slides:
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"):
                                                file_text += shape.text + "\n"

                        except:
                            pass  

           
            elif uploaded_file.type.startswith("image"):
                image = Image.open(uploaded_file)
                st.image(image)
            

              

        user_input = st.chat_input("Ask something...")

        if user_input or uploaded_file:

            
            if image:

                    if not user_input:
                        st.warning("⚠️ Please ask something about the image.")
                        st.stop()

                    prompt = user_input

                    st.chat_message("user").markdown(prompt)

                    inputs = processor(
                        image,
                        text=prompt,
                        return_tensors="pt"
                    )

                    output = model.generate(
                        **inputs,
                        max_length=100,
                        num_beams=5
                    )

                    caption = processor.decode(output[0], skip_special_tokens=True)

                    full_prompt = f"""
                User question: {prompt}

                Image description: {caption}

                Now explain this image in a detailed, friendly, human-like way.
                Talk naturally like a helpful assistant. Give more details.
                """

                    response = client.chat.completions.create(
                        model="llama-3.1-8b-instant",
                        messages=[{"role": "user", "content": full_prompt}]
                    )

                    reply = response.choices[0].message.content

                    st.chat_message("assistant").markdown(reply)

                    st.session_state.messages.append({"role": "user", "content": prompt})
                    st.session_state.messages.append({"role": "assistant", "content": reply})

            else:
                combined = user_input or ""

                if file_text:
                    combined += f"\n\nFile Content:\n{file_text[:2000]}"

                st.session_state.messages.append({"role": "user", "content": combined})
                st.chat_message("user").markdown(combined)

                response = client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=st.session_state.messages
                )

                reply = response.choices[0].message.content

                st.chat_message("assistant").markdown(reply)
                st.session_state.messages.append({"role": "assistant", "content": reply})
    
elif selected_app == "Voice Assistant":
    st.title("Start With Voice Assistant")
    r= sr.Recognizer()
    def get_voice_input():
            with sr.Microphone() as source:
                st.info("🎤 Listening... Please speak.")
                audio = r.listen(source)
                try:
                    return r.recognize_google(audio, language="en-in")
                except sr.UnknownValueError:
                    return "Sorry, I didn't catch that."
                
    def process_input(prompt):
        with st.spinner("💡 Groq is thinking..."):

            response = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )

            return response.choices[0].message.content
    def speak(text):
            sentences = text.split('. ')
            for sentence in sentences:
                if sentence.strip() == "":
                    continue
                tts = gTTS(text=sentence, lang='en')
                with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
                    tts.save(fp.name)
                    st.audio(fp.name, format="audio/mp3")
                    
            
            
    if st.button("🎤 Talk to Gemini"):
            user_input = get_voice_input()
            if user_input.strip() != "":
                response_text = process_input(user_input)
                speak(response_text)
                st.markdown("## 🤖 Gemini's Response:")
                st.success(response_text)
        
